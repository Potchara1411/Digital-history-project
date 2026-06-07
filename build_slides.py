"""Generate the presentation deck (ai_discourse_slides.pptx) from the project's
charts. Run after ai_final.py so the figures/ panels are up to date:

    python ai_final.py
    python build_slides.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = Path(__file__).resolve().parent
FIG = BASE / "figures"
OUT = BASE / "ai_discourse_slides.pptx"

# Palette (matches the website / charts)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1D, 0x24)
MUTED = RGBColor(0x5B, 0x64, 0x70)
GOLD = RGBColor(0xB8, 0x86, 0x0B)
BLUE = RGBColor(0x2E, 0x75, 0xB6)
RED = RGBColor(0xE7, 0x4C, 0x3C)
PURPLE = RGBColor(0x7E, 0x4F, 0xA8)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
FONT = "Calibri"

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def bg(slide, color=WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def accent_bar(slide, x, y, w, h, color=GOLD):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def text(slide, s, x, y, w, h, size=18, color=INK, bold=False, italic=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0, font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(s.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        run = p.add_run()
        run.text = line
        f = run.font
        f.size, f.color.rgb, f.bold, f.italic, f.name = Pt(size), color, bold, italic, font
    return tb


def bullets(slide, items, x, y, w, h, size=18, color=INK, gap=8):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        txt, c = item if isinstance(item, tuple) else (item, color)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        dot = p.add_run()
        dot.text = "•  "
        dot.font.size, dot.font.color.rgb, dot.font.bold, dot.font.name = Pt(size), GOLD, True, FONT
        run = p.add_run()
        run.text = txt
        run.font.size, run.font.color.rgb, run.font.name = Pt(size), c or color, FONT
    return tb


def picture_centered(slide, path, top, max_w=Inches(11.0), max_h=Inches(4.35)):
    from PIL import Image
    iw, ih = Image.open(path).size
    scale = min(max_w / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    slide.shapes.add_picture(str(path), int((EMU_W - w) / 2), top, width=w, height=h)


def header(slide, label, title, title_size=28):
    bg(slide)
    accent_bar(slide, Inches(0.6), Inches(0.55), Inches(0.16), Inches(0.62))
    text(slide, label, Inches(0.95), Inches(0.5), Inches(11.6), Inches(0.4),
         size=14, color=GOLD, bold=True)
    text(slide, title, Inches(0.95), Inches(0.82), Inches(11.9), Inches(0.8),
         size=title_size, color=INK, bold=True)


def chart_slide(label, title, panel, points):
    s = prs.slides.add_slide(BLANK)
    header(s, label, title)
    bullets(s, points, Inches(0.95), Inches(1.72), Inches(11.6), Inches(1.0), size=16)
    picture_centered(s, panel, Inches(2.85))
    return s


def text_slide(label, title, items, size=20, gap=16, coda=None):
    s = prs.slides.add_slide(BLANK)
    header(s, label, title, title_size=30)
    bullets(s, items, Inches(1.15), Inches(2.0), Inches(11.3), Inches(4.3), size=size, gap=gap)
    if coda:
        text(s, coda, Inches(1.15), Inches(6.05), Inches(11), Inches(0.8),
             size=18, color=GOLD, italic=True)
    return s


# ---------- 1 · Title ----------
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s, Inches(0.0), Inches(0.0), EMU_W, Inches(0.22))
accent_bar(s, Inches(0.0), Inches(7.28), EMU_W, Inches(0.22))
text(s, "KAIST · DIGITAL HISTORY", Inches(1), Inches(1.95), Inches(11.3), Inches(0.4),
     size=15, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
text(s, "From Scientific Optimism\nto Geopolitical Anxiety", Inches(1), Inches(2.45),
     Inches(11.3), Inches(2.0), size=44, color=INK, bold=True, align=PP_ALIGN.CENTER)
text(s, "A Computational History of AI Discourse  ·  1990–2023", Inches(1), Inches(4.5),
     Inches(11.3), Inches(0.6), size=21, color=MUTED, align=PP_ALIGN.CENTER)
text(s, "potchara1411.github.io/Digital-history-project", Inches(1), Inches(5.45),
     Inches(11.3), Inches(0.4), size=15, color=BLUE, align=PP_ALIGN.CENTER)

# ---------- 2 · The gap ----------
text_slide("THE GAP", "Filling a gap in Crawford's Atlas of AI", [
    ("Kate Crawford's Atlas of AI (2021) argues AI is a system of power and extraction "
     "— but qualitatively. It never proves computationally WHEN public discourse shifted, "
     "or HOW FAST.", INK),
    ("This project supplies that evidence: 30 years of book data + 13 years of global "
     "news data, locating the turning points and measuring their speed.", INK),
    ("It flips the question — not “what can AI do for historians?” but “what can historians "
     "do with AI discourse as a primary source?”", INK),
    ("Secondary frame: AI anxiety (2010s–20s) reuses the template of nuclear anxiety "
     "(1950s–60s).", MUTED),
], size=19)

# ---------- 3 · Data & method ----------
text_slide("DATA & METHOD", "Two kinds of evidence", [
    ("Conventional sources frame the project: Crawford's Atlas of AI; print archives "
     "(Wired, TIME, The Economist); government & policy documents.", INK),
    ("Digital sources provide the numbers: Google Ngram (books, 1990–2019); Media Cloud "
     "(news, 2010–2023); GDELT (normalized coverage + tone).", INK),
    ("Method: quantitative pattern detection + qualitative interpretation — adapting "
     "Dr. Woo's semantic-shift approach to Cold War newspapers.", INK),
    ("Fully reproducible: every chart is built by one Python script from public CSVs.", MUTED),
], size=19)

# ---------- 4–6 · the three acts ----------
chart_slide("ACT I · 1990–2013", "The Winter and the Rebranding",
            FIG / "panel1_ngram.png", [
                ("“Artificial intelligence” DECLINES through the 1990s AI Winter", BLUE),
                ("“Machine learning” rises and overtakes it ~2013 — the field rebrands around "
                 "concrete method", RED),
            ])
chart_slide("ACT II · 2012–2021", "Back into the Light",
            FIG / "panel2_volume.png", [
                ("Deep learning works (2012 ImageNet); AlphaGo stuns the public (2016)", INK),
                ("News coverage climbs from a few thousand to tens of thousands of articles a year",
                 INK),
            ])
chart_slide("ACT III · 2022–2023", "The ChatGPT Explosion",
            FIG / "panel3_share.png", [
                ("ChatGPT (Nov 2022) puts AI in everyone’s hands", INK),
                ("AI’s SHARE of all news more than doubles in a year — to ~1% of everything in "
                 "the press", PURPLE),
            ])

# ---------- 7 · The twist: tone ----------
chart_slide("THE TWIST", "It was never optimism",
            FIG / "panel5_tone.png", [
                ("Average tone of AI coverage was NEGATIVE every single year — even in 2010", RED),
                ("There was no golden age of optimism to fall from — the “optimism → fear” story "
                 "is wrong", INK),
            ])

# ---------- 8 · The argument ----------
text_slide("THE ARGUMENT", "Scale and salience — not sentiment", [
    ("AI discourse did not drift from hope to fear; it was wary from the start.", INK),
    ("What changed was SCALE — attention exploded after AlphaGo (2016) and ChatGPT (2022–23).",
     INK),
    ("And POLITICAL SALIENCE — AI shifted from a scientific keyword (Deep Blue, research) to a "
     "geopolitical one (US–China, existential risk, regulation).", INK),
    ("The anxiety isn’t new; its scale and its politics are — the same template as 1950s–60s "
     "nuclear fear.", GOLD),
], size=19, gap=18)

# ---------- 9 · Robustness ----------
chart_slide("ROBUSTNESS CHECK", "Does a second source agree?",
            FIG / "panel4_robustness.png", [
                ("GDELT (normalized) and Media Cloud track each other across 2017–2023", INK),
                ("In 2023 they land at nearly the same value — ~1.04% vs ~1.05% — so the surge is "
                 "real, not a single-source artifact", INK),
            ])

# ---------- 10 · Why not GDELT raw ----------
text_slide("A METHODOLOGICAL CHOICE", "Why we did NOT trust GDELT’s raw counts", [
    ("GDELT’s AI event count jumps 347 (2013) → 52,830 (2014) → 3.9 MILLION (2015) — a "
     "~10,000× leap in two years.", INK),
    ("That is not AI exploding; it is GDELT 2.0 expanding its own sources in 2015. The "
     "instrument changed, not the world.", RED),
    ("Fix: use only its time-comparable metrics — normalized share and average tone — and flag "
     "the small early samples.", INK),
    ("Source criticism is the scholarship: a measurement is only as good as the instrument’s "
     "consistency over time.", MUTED),
], size=19)

# ---------- 11 · Conclusion ----------
text_slide("CONCLUSION", "AI discourse as a primary source", [
    ("Books — slow, deliberate — recorded a quiet rebranding: AI → machine learning → AI.", INK),
    ("News — fast, reactive — recorded an explosion in scale and political salience after ChatGPT.",
     INK),
    ("The shift was scale and politics, not optimism to fear — extending Crawford’s argument with "
     "quantitative evidence.", INK),
], size=20, gap=18,
    coda="Treating AI coverage as a primary source shows how societies negotiate power, fear, "
         "and technological change.")

# ---------- 12 · Thank you ----------
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s, Inches(0.0), Inches(0.0), EMU_W, Inches(0.22))
accent_bar(s, Inches(0.0), Inches(7.28), EMU_W, Inches(0.22))
text(s, "Thank you", Inches(1), Inches(2.5), Inches(11.3), Inches(1.0),
     size=44, color=INK, bold=True, align=PP_ALIGN.CENTER)
text(s, "Explore the interactive exhibit", Inches(1), Inches(3.8), Inches(11.3), Inches(0.5),
     size=20, color=MUTED, align=PP_ALIGN.CENTER)
text(s, "potchara1411.github.io/Digital-history-project", Inches(1), Inches(4.3),
     Inches(11.3), Inches(0.5), size=20, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
text(s, "Code & data: github.com/Potchara1411/Digital-history-project", Inches(1),
     Inches(4.95), Inches(11.3), Inches(0.4), size=15, color=MUTED, align=PP_ALIGN.CENTER)

# ---------- Presenter notes (same script as PRESENTATION_SCRIPT.md) ----------
NOTES = [
    "My project is a computational history of artificial intelligence — not the technology, "
    "but the discourse. I treat media coverage of AI as a primary source, and ask when and how "
    "the way we talk about AI actually changed, across more than thirty years.",

    "My starting point is Kate Crawford's Atlas of AI. She argues powerfully that AI is a system "
    "of power and extraction — but it's a qualitative argument. She never shows, computationally, "
    "WHEN public discourse shifted or HOW FAST. That's the gap I fill, with 30 years of book data "
    "and 13 years of global news data. And I flip the usual question: not what AI can do for "
    "historians, but what historians can do with AI discourse as a source. My second frame: this "
    "looks a lot like nuclear anxiety in the 1950s and 60s.",

    "Two kinds of evidence. Conventional sources frame the project — Crawford's book, print "
    "archives like Wired, TIME and the Economist, and policy documents. Digital sources give me "
    "the numbers — Google Ngram for books, Media Cloud for news, and GDELT. The method pairs "
    "quantitative pattern detection with qualitative interpretation, adapting Dr. Woo's work on "
    "semantic shifts in Cold War newspapers. And it's all reproducible — one script, public data.",

    "Act one — the first surprise: the story starts with DECLINE. The blue line, 'artificial "
    "intelligence,' falls through the 1990s — the AI Winter, when the term became almost "
    "embarrassing. But the red line, 'machine learning,' rises and overtakes it around 2013. "
    "The field didn't quit; it rebranded around concrete method.",

    "Act two: how AI came back. The technology started working — deep learning around 2012 — and "
    "there was public spectacle, above all AlphaGo beating the world's best Go player in 2016. "
    "News coverage climbs from a few thousand to tens of thousands of articles a year.",

    "Act three: the explosion. ChatGPT, late 2022. But raw counts can mislead, so I measure AI's "
    "SHARE of all coverage. Even controlling for the growth of news, it more than doubled in one "
    "year — to about 1% of everything in the press. A real shift in attention.",

    "Now the twist — and this is where I push back on the obvious story. The conventional narrative "
    "is optimism turning to fear. But look at the tone of AI coverage: it's negative EVERY single "
    "year, around minus five, even back in 2010. There was no golden age of optimism to fall from. "
    "I'm honest about the data — the early years are small samples, marked hollow — but the pattern "
    "is clear.",

    "So here's my argument. AI discourse didn't drift from hope to fear; it was wary from the start. "
    "What actually changed was SCALE — the explosion of attention after AlphaGo and ChatGPT — and "
    "POLITICAL SALIENCE: AI went from a scientific keyword to a geopolitical one, tied to US-China "
    "competition and existential risk. The anxiety isn't new; its scale and its politics are. And "
    "that's the same template societies used for nuclear fear in the 1950s and 60s.",

    "A skeptic should ask: what if the 2023 spike is just a Media Cloud quirk? So I test it against "
    "GDELT's normalized measure. Two datasets, different teams, different methods — yet they track "
    "closely and in 2023 land at almost the same value, 1.04 versus 1.05 percent. The surge is real.",

    "The slide I'm proudest of. I almost used GDELT as my main news source — but its raw counts go "
    "from 347 to 53,000 to 3.9 MILLION in two years. That's not AI exploding; it's GDELT expanding "
    "its own sources in 2015. The instrument changed, not the world. So I use only its "
    "time-comparable metrics. This source criticism is the actual scholarship.",

    "To conclude: books recorded a quiet rebranding; news recorded an explosion in scale and "
    "political salience. The shift was scale and politics, not optimism to fear. By treating AI "
    "coverage as a primary source, I extend Crawford's argument with the quantitative evidence it "
    "was missing — and show how societies negotiate power, fear, and technological change.",

    "That's my project. The interactive exhibit and all the code and data are online. Thank you — "
    "happy to take questions.",
]
for slide, note in zip(prs.slides, NOTES):
    slide.notes_slide.notes_text_frame.text = note

prs.save(OUT)
print(f"Saved {OUT.name}  ({len(prs.slides._sldIdLst)} slides, notes embedded)")
