"""Build the 4-minute in-progress presentation (ai_discourse_4min.pptx) — a tight
5-slide promo version, separate from the full deck (build_slides.py).

    python ai_final.py
    python build_slides_4min.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = Path(__file__).resolve().parent
FIG = BASE / "figures"
OUT = BASE / "ai_discourse_4min.pptx"

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1D, 0x24)
MUTED = RGBColor(0x5B, 0x64, 0x70)
GOLD = RGBColor(0xB8, 0x86, 0x0B)
BLUE = RGBColor(0x2E, 0x75, 0xB6)
RED = RGBColor(0xE7, 0x4C, 0x3C)
PURPLE = RGBColor(0x7E, 0x4F, 0xA8)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
FONT = "Calibri"

prs = Presentation()
prs.slide_width, prs.slide_height = EMU_W, EMU_H
BLANK = prs.slide_layouts[6]


def bg(slide, color=WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def accent_bar(slide, x, y, w, h, color=GOLD):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp


def text(slide, s, x, y, w, h, size=18, color=INK, bold=False, italic=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, line in enumerate(s.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing
        r = p.add_run(); r.text = line
        r.font.size, r.font.color.rgb, r.font.bold, r.font.italic, r.font.name = (
            Pt(size), color, bold, italic, FONT)
    return tb


def bullets(slide, items, x, y, w, h, size=20, gap=12):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, (txt, c) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.05
        d = p.add_run(); d.text = "•  "
        d.font.size, d.font.color.rgb, d.font.bold, d.font.name = Pt(size), GOLD, True, FONT
        r = p.add_run(); r.text = txt
        r.font.size, r.font.color.rgb, r.font.name = Pt(size), c, FONT
    return tb


def picture_centered(slide, path, top, max_w=Inches(10.4), max_h=Inches(4.4)):
    from PIL import Image
    iw, ih = Image.open(path).size
    scale = min(max_w / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    slide.shapes.add_picture(str(path), int((EMU_W - w) / 2), top, width=w, height=h)


def header(slide, label, title):
    bg(slide)
    accent_bar(slide, Inches(0.6), Inches(0.55), Inches(0.16), Inches(0.62))
    text(slide, label, Inches(0.95), Inches(0.5), Inches(11.6), Inches(0.4), size=14, color=GOLD, bold=True)
    text(slide, title, Inches(0.95), Inches(0.82), Inches(11.9), Inches(0.8), size=30, color=INK, bold=True)


# ---- 1 · Title ----
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s, 0, 0, EMU_W, Inches(0.22)); accent_bar(s, 0, Inches(7.28), EMU_W, Inches(0.22))
text(s, "KAIST · DIGITAL HISTORY  ·  PROJECT IN PROGRESS", Inches(1), Inches(1.95), Inches(11.3),
     Inches(0.4), size=15, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
text(s, "From Scientific Optimism\nto Geopolitical Anxiety", Inches(1), Inches(2.45), Inches(11.3),
     Inches(2.0), size=44, color=INK, bold=True, align=PP_ALIGN.CENTER)
text(s, "A Computational History of AI Discourse  ·  1990–2023", Inches(1), Inches(4.5), Inches(11.3),
     Inches(0.6), size=21, color=MUTED, align=PP_ALIGN.CENTER)
text(s, "potchara1411.github.io/Digital-history-project", Inches(1), Inches(5.45), Inches(11.3),
     Inches(0.4), size=15, color=BLUE, align=PP_ALIGN.CENTER)

# ---- 2 · What's new ----
s = prs.slides.add_slide(BLANK)
header(s, "WHAT'S NEW", "Filling a gap in Crawford's Atlas of AI")
bullets(s, [
    ("Kate Crawford argues AI is a system of power — but qualitatively. She never shows "
     "WHEN public discourse shifted, or how fast.", INK),
    ("I add the quantitative timeline: 30 years of books (Ngram) + 13 years of global news "
     "(Media Cloud, GDELT).", INK),
    ("And I flip the question — not “what can AI do for historians?” but “what can historians "
     "do with AI discourse as a primary source?”", INK),
], Inches(1.15), Inches(2.05), Inches(11.3), Inches(4), size=21, gap=20)

# ---- 3 · The rebrand ----
s = prs.slides.add_slide(BLANK)
header(s, "FINDING 1", "AI didn't just grow — it declined, then rebranded")
bullets(s, [
    ("“Artificial intelligence” DECLINED in the 1990s — the AI Winter", BLUE),
    ("“Machine learning” rose and overtook it by 2013 — the field quietly rebranded around "
     "concrete method", RED),
], Inches(0.95), Inches(1.72), Inches(11.6), Inches(1.0), size=16)
picture_centered(s, FIG / "panel1_ngram.png", Inches(2.9))

# ---- 4 · The explosion ----
s = prs.slides.add_slide(BLANK)
header(s, "FINDING 2", "Then the explosion — and it's real")
bullets(s, [
    ("After AlphaGo (2016) and ChatGPT (2022), AI’s SHARE of all news more than doubled — to "
     "~1% in 2023", PURPLE),
    ("I control for total news growth, so this is a genuine shift in attention, not just "
     "“more news”", INK),
], Inches(0.95), Inches(1.72), Inches(11.6), Inches(1.0), size=16)
picture_centered(s, FIG / "panel3_share.png", Inches(2.9))

# ---- 5 · The twist ----
s = prs.slides.add_slide(BLANK)
header(s, "FINDING 3 · THE TWIST", "It was never optimism")
bullets(s, [
    ("Tone of AI coverage was NEGATIVE every single year — even in 2010. No optimism to lose.",
     RED),
    ("So the shift wasn’t hope → fear. It was SCALE + POLITICAL SALIENCE: AI went from a "
     "scientific keyword to a geopolitical one.", INK),
], Inches(0.95), Inches(1.72), Inches(11.6), Inches(1.0), size=16)
picture_centered(s, FIG / "panel5_tone.png", Inches(2.9))

# ---- 6 · Method credibility ----
s = prs.slides.add_slide(BLANK)
header(s, "WHY IT'S RIGOROUS", "I rejected a source — on purpose")
bullets(s, [
    ("I almost used GDELT — but its AI counts jump 347 → 53,000 → 3.9 MILLION in just two years",
     INK),
    ("That’s the database expanding its sources in 2015 — the instrument changed, not the world",
     RED),
    ("Catching that — source criticism — is the real scholarship behind the project", MUTED),
], Inches(1.15), Inches(2.05), Inches(11.3), Inches(4), size=21, gap=20)

# ---- 5 · Close ----
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s, 0, 0, EMU_W, Inches(0.22)); accent_bar(s, 0, Inches(7.28), EMU_W, Inches(0.22))
text(s, "See the full exhibit — and tell me where it's weak", Inches(1), Inches(2.2), Inches(11.3),
     Inches(1.0), size=30, color=INK, bold=True, align=PP_ALIGN.CENTER)
text(s, "Five interactive charts · a GDELT robustness check · full source criticism", Inches(1),
     Inches(3.5), Inches(11.3), Inches(0.5), size=18, color=MUTED, align=PP_ALIGN.CENTER)
text(s, "potchara1411.github.io/Digital-history-project", Inches(1), Inches(4.2), Inches(11.3),
     Inches(0.5), size=22, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
text(s, "github.com/Potchara1411/Digital-history-project", Inches(1), Inches(4.95), Inches(11.3),
     Inches(0.4), size=14, color=MUTED, align=PP_ALIGN.CENTER)

NOTES = [
    "Hi — my DH project asks a simple question: when did 'AI' become inescapable, and was it ever "
    "the hopeful story we assume? I trace how we've talked about AI for thirty years — in books and "
    "in the news. [~20s]",

    "My hook is Kate Crawford's Atlas of AI. She argues AI is a system of power — but it's a "
    "qualitative argument; she never shows computationally WHEN the discourse shifted or how fast. "
    "That's my gap: 30 years of book data plus 13 years of global news data. And I flip the usual "
    "question — not what AI can do for historians, but what historians can do with AI discourse as a "
    "primary source. [~45s]",

    "Finding one, and it's a surprise: AI didn't just grow. It DECLINED in the 1990s — the AI Winter "
    "— and 'machine learning' overtook it by 2013, as the field quietly rebranded around concrete "
    "method. [~40s]",

    "Finding two: then the explosion. After AlphaGo and ChatGPT, AI's SHARE of all news more than "
    "doubled — to about 1% in 2023. And I control for the total growth of news, so this is a real "
    "shift in attention, not just more articles being published. [~40s]",

    "Finding three — the twist I'm proudest of. The conventional story is optimism turning to fear. "
    "But the tone of AI coverage was negative every single year, even back in 2010. There was no "
    "golden age of optimism to lose. So the real shift wasn't sentiment — it was SCALE and POLITICAL "
    "SALIENCE: AI went from a scientific keyword to a geopolitical one, tied to US-China competition "
    "and existential risk. [~50s]",

    "One methodological point I'm proud of. I almost used GDELT as my main news source — but its AI "
    "counts jump from 347, to 53,000, to 3.9 MILLION in two years. That's not AI exploding; it's the "
    "database expanding its own sources in 2015. The instrument changed, not the world. Catching "
    "that — source criticism — is the real scholarship. [~35s]",

    "It's all a live web exhibit — with a robustness check against a second dataset and full source "
    "criticism. Please go explore it, and tell me where the argument is weak — I'd genuinely love "
    "your pushback. Thank you. [~25s]",
]
for slide, note in zip(prs.slides, NOTES):
    slide.notes_slide.notes_text_frame.text = note

prs.save(OUT)
print(f"Saved {OUT.name}  ({len(prs.slides._sldIdLst)} slides)")
